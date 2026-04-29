"""Generate the Phase-1 status presentation from the BAM template.

Run from the project root with the project venv:
    python scripts/build_status_slides.py
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

SRC = Path("docs/BAM_Vorlage.pptx")
DST = Path("docs/Phase1_Status_2026-04.pptx")

FOOTER_TEXT = "Phase 1 Status  ·  Technical-PDF Document Analysis"
DATE_TEXT = "April 2026"

# --- Palette ------------------------------------------------------------
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
BAM_RED = RGBColor(0xC6, 0x1F, 0x37)
LIGHT_BG = RGBColor(0xE7, 0xEC, 0xF1)
ACCENT_LIGHT = RGBColor(0xFC, 0xE3, 0xE8)
MID_GRAY = RGBColor(0x9A, 0xA4, 0xB1)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# --- Helpers ------------------------------------------------------------
def get_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def remove_placeholder(slide, idx):
    ph = get_placeholder(slide, idx)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def clone_layout_placeholder(slide, layout_ph_idx):
    """Copy a placeholder element from the slide's layout onto the slide.

    Used for footer/date/slide-number placeholders that python-pptx does not
    auto-clone when you call add_slide().
    """
    layout = slide.slide_layout
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == layout_ph_idx:
            cloned = deepcopy(ph._element)
            slide.shapes._spTree.append(cloned)
            return slide.placeholders[layout_ph_idx]
    return None


def add_footer_block(slide, *, footer_text=FOOTER_TEXT, date_text=DATE_TEXT):
    """Add the BAM-style date / footer / slide-number trio to a slide."""
    # Date
    date_ph = clone_layout_placeholder(slide, 10)
    if date_ph is not None:
        date_ph.text_frame.text = date_text
        for p in date_ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)
                r.font.color.rgb = MID_GRAY
    # Footer
    foot_ph = clone_layout_placeholder(slide, 11)
    if foot_ph is not None:
        foot_ph.text_frame.text = footer_text
        for p in foot_ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)
                r.font.color.rgb = MID_GRAY
    # Slide number — keep the layout's existing field; just clone
    clone_layout_placeholder(slide, 12)


def set_speaker_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


def style_paragraph(p, *, size=18, bold=False, italic=False, color=TEXT_DARK, align=None):
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def write_paragraphs(text_frame, paragraphs, *, size=16, color=TEXT_DARK,
                     align=None, bullet=False, space_after=6):
    text_frame.clear()
    text_frame.word_wrap = True
    for i, item in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = ("• " + item) if bullet else item
        p.space_after = Pt(space_after)
        style_paragraph(p, size=size, color=color, align=align)


def fill_box(shape, fill_color, line_color=DARK_BLUE, line_width_pt=0.75):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width_pt)


def labeled_box(slide, x_cm, y_cm, w_cm, h_cm, text, *,
                fill=LIGHT_BG, line=DARK_BLUE, fg=TEXT_DARK,
                size=14, bold=False, shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    fill_box(box, fill, line_color=line)
    tf = box.text_frame
    tf.text = text
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    style_paragraph(tf.paragraphs[0], size=size, bold=bold,
                    color=fg, align=PP_ALIGN.CENTER)
    return box


def set_title(slide, text, *, size=28, color=DARK_BLUE):
    if slide.shapes.title is None:
        return
    tf = slide.shapes.title.text_frame
    tf.text = text
    style_paragraph(tf.paragraphs[0], size=size, bold=True, color=color)


# --- Build --------------------------------------------------------------
prs = Presentation(SRC)

# ---------- Slide 1: Title (edit existing) ----------
title_slide = prs.slides[0]
for shape in title_slide.shapes:
    if not shape.has_text_frame:
        continue
    name = shape.name
    if name.startswith("Titel"):
        shape.text_frame.text = "Technical-PDF Document Analysis — Phase 1"
        style_paragraph(shape.text_frame.paragraphs[0], size=30,
                        bold=True, color=DARK_BLUE)
    elif name.startswith("Untertitel"):
        shape.text_frame.text = (
            "Structured extraction as the foundation for claim verification"
        )
        style_paragraph(shape.text_frame.paragraphs[0], size=18, color=TEXT_DARK)
    elif name == "Textplatzhalter 13":
        shape.text_frame.text = "April 2026  ·  Mohamed Tababi"
        style_paragraph(shape.text_frame.paragraphs[0], size=14, color=TEXT_DARK)
    # the FB 3.3 free-standing text box stays untouched


# ---------- Slide 2: Why & Goal ----------
s = prs.slides.add_slide(prs.slide_layouts[6])  # 1-spaltig universal
set_title(s, "Why & Goal")
body = get_placeholder(s, 13)
write_paragraphs(
    body.text_frame,
    [
        "Technical safety PDFs mix text, tables, formulas, figures, "
        "and technical drawings — today not machine-usable.",
        "Phase 1 goal: extract every modality into one stable, "
        "structured output.",
        "Long-term goal: verify technical claims and trace each one "
        "back to its supporting evidence — text and visual.",
        "Driver: support BAM's approval of safety-critical "
        "transport containers.",
    ],
    size=18, bullet=True, space_after=10,
)


# ---------- Slide 3: System Architecture ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "System Architecture — 4 Independent Layers")
remove_placeholder(s, 13)

# Stack diagram on the left half
stack_x = 1.5
stack_y = 3.2
stack_w = 9.0
stack_h = 1.4
gap = 0.25
layers = [
    ("Layer 4  ·  Agent / Routing", LIGHT_BG, TEXT_DARK, False),
    ("Layer 3  ·  Storage", LIGHT_BG, TEXT_DARK, False),
    ("Layer 2  ·  Embedding", LIGHT_BG, TEXT_DARK, False),
    ("Layer 1  ·  Extraction      (current scope)", BAM_RED, WHITE, True),
]
for i, (label, fill, fg, bold) in enumerate(layers):
    labeled_box(
        s, stack_x, stack_y + i * (stack_h + gap), stack_w, stack_h,
        label, fill=fill, fg=fg, size=15, bold=bold,
    )

# Right-side notes
note_box = s.shapes.add_textbox(Cm(11.5), Cm(3.2), Cm(13.0), Cm(9.0))
write_paragraphs(
    note_box.text_frame,
    [
        "Each layer is independently replaceable.",
        "Later layers depend on the Phase-1 output contract, "
        "not on extraction internals.",
        "Phase 1 produces stable JSON sidecars the rest of the "
        "system can build on.",
        "Tools inside Layer 1 are also swappable — see the "
        "next slides.",
    ],
    size=15, bullet=True, space_after=10,
)


# ---------- Slide 4: Extraction Pipeline ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Extraction Pipeline — 4 Staged Commands")
remove_placeholder(s, 13)

stages = [
    ("segment", "render pages,\nrun layout segmenter"),
    ("extract-text", "text · tables · formulas\nvia configured extractors"),
    ("describe-figures", "figures · diagrams · drawings\nvia VLM"),
    ("assemble", "deterministic merge\n→ content_list.json"),
]
n = len(stages)
total_w = 22.5
gap_cm = 0.45
box_w = (total_w - gap_cm * (n - 1)) / n
box_h = 3.2
top = 4.0
left = 1.5

for i, (head, sub) in enumerate(stages):
    x = left + i * (box_w + gap_cm)
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(x), Cm(top), Cm(box_w), Cm(box_h),
    )
    fill_box(box, LIGHT_BG)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = head
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT_DARK

    if i < n - 1:
        ax = x + box_w + 0.05
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Cm(ax), Cm(top + box_h / 2 - 0.4),
            Cm(gap_cm - 0.1), Cm(0.8),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# Footnote bullets
foot = s.shapes.add_textbox(Cm(1.5), Cm(8.2), Cm(22.5), Cm(4.0))
write_paragraphs(
    foot.text_frame,
    [
        "Each stage runs in its own process — GPU memory is "
        "released between stages.",
        "6 swappable roles (renderer, segmenter, text/table/formula "
        "extractor, figure descriptor) — chosen by YAML config, "
        "not by heuristic.",
        "Defaults today: pymupdf · MinerU 3.1 · MinerU passthrough "
        "for text/tables/formulas · Qwen2.5-VL for figures.",
    ],
    size=13, bullet=True, space_after=6,
)


# ---------- Slide 5: Code Architecture ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Code Architecture — Three Concentric Layers")
remove_placeholder(s, 13)

# Top: CLI band spanning all three columns
cli = labeled_box(
    s, 1.7, 2.5, 22.0, 0.9,
    "python -m extraction  segment | extract-text | describe-figures | assemble",
    fill=DARK_BLUE, fg=WHITE, line=DARK_BLUE, size=12, bold=True,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
)

# Three columns
col_specs = [
    ("Core",
     "configuration · contracts",
     [
         "interfaces.py     — 6 Protocols (roles)",
         "registry.py       — name → adapter dispatch",
         "config.py         — YAML → ExtractionConfig",
         "models.py         — Region · Element · ContentList",
         "output.py         — sidecar writer + crop scaling",
     ]),
    ("Stages",
     "orchestration",
     [
         "segment.py            — render + layout",
         "extract_text.py       — text · tables · formulas",
         "describe_figures.py   — figures · diagrams · drawings",
         "assemble.py           — deterministic merge",
     ]),
    ("Adapters",
     "concrete tools (lazy-imported)",
     [
         "pymupdf_renderer",
         "mineru25_segmenter",
         "pymupdf_text_segmenter",
         "olmocr2_text",
         "qwen25vl_figure",
         "stubs (noop)",
     ]),
]
col_w = 7.20
col_gap = 0.40
col_x0 = 1.7
col_y = 3.7
col_h = 7.5

for i, (head, sub, items) in enumerate(col_specs):
    x = col_x0 + i * (col_w + col_gap)
    # Container
    cont = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Cm(x), Cm(col_y), Cm(col_w), Cm(col_h))
    fill_box(cont, LIGHT_BG)
    cont.text_frame.text = ""  # we'll write into separate textboxes overlaid

    # Header text on top of container
    head_tb = s.shapes.add_textbox(Cm(x), Cm(col_y + 0.2),
                                   Cm(col_w), Cm(0.8))
    head_tb.text_frame.word_wrap = True
    head_tb.text_frame.text = head
    style_paragraph(head_tb.text_frame.paragraphs[0],
                    size=15, bold=True, color=DARK_BLUE,
                    align=PP_ALIGN.CENTER)
    # Subhead
    sub_p = head_tb.text_frame.add_paragraph()
    sub_p.text = sub
    style_paragraph(sub_p, size=10, italic=True,
                    color=MID_GRAY, align=PP_ALIGN.CENTER)

    # File list
    list_tb = s.shapes.add_textbox(Cm(x + 0.4), Cm(col_y + 1.7),
                                   Cm(col_w - 0.8), Cm(col_h - 2.0))
    list_tb.text_frame.word_wrap = True
    for j, item in enumerate(items):
        p = list_tb.text_frame.paragraphs[0] if j == 0 else list_tb.text_frame.add_paragraph()
        p.text = item
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(10)
            r.font.color.rgb = TEXT_DARK
            r.font.name = "Consolas"

# Bottom note
bot = s.shapes.add_textbox(Cm(1.7), Cm(11.4), Cm(22.0), Cm(0.8))
bot.text_frame.text = (
    "+ tests/  (14 modules, GPU paths marker-gated)"
)
style_paragraph(bot.text_frame.paragraphs[0],
                size=11, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)


# ---------- Slide 6: Modularity Guarantee ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Modularity — How Tool Swaps Stay Safe")
remove_placeholder(s, 13)

# 4 step boxes across the slide, each with a code snippet
step_specs = [
    ("1.  Protocol",
     "interfaces.py",
     "class Segmenter(Protocol):\n"
     "    @property\n"
     "    def tool_name(self) -> str: ...\n"
     "    def segment(self, pdf): ..."),
    ("2.  Adapter",
     "adapters/mineru25_segmenter.py",
     "@register_segmenter(\"mineru25\")\n"
     "class MinerU25Segmenter:\n"
     "    tool_name = \"mineru25\"\n"
     "    def segment(self, pdf): ..."),
    ("3.  Config",
     "extraction_config.yaml",
     "extraction:\n"
     "  segmenter: mineru25\n"
     "  text_extractor: olmocr2\n"
     "  figure_descriptor: qwen25vl"),
    ("4.  Pipeline",
     "stages/segment.py",
     "seg = get_segmenter(\n"
     "    cfg.segmenter,\n"
     "    **cfg.adapter_kwargs,\n"
     ")\nregions = seg.segment(pdf)"),
]

n_steps = len(step_specs)
total_w = 22.5
gap_w = 0.30
step_w = (total_w - gap_w * (n_steps - 1)) / n_steps   # ~5.4 cm
step_h = 5.5
step_top = 3.3
step_left = 1.5

step_fills = [LIGHT_BG, LIGHT_BG, LIGHT_BG, ACCENT_LIGHT]

for i, (head, where, code) in enumerate(step_specs):
    x = step_left + i * (step_w + gap_w)
    # Container
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Cm(x), Cm(step_top), Cm(step_w), Cm(step_h))
    fill_box(box, step_fills[i])

    # Title
    title_tb = s.shapes.add_textbox(Cm(x + 0.2), Cm(step_top + 0.15),
                                    Cm(step_w - 0.4), Cm(0.7))
    title_tb.text_frame.text = head
    style_paragraph(title_tb.text_frame.paragraphs[0],
                    size=13, bold=True, color=DARK_BLUE)

    # File path
    where_tb = s.shapes.add_textbox(Cm(x + 0.2), Cm(step_top + 0.85),
                                    Cm(step_w - 0.4), Cm(0.5))
    where_tb.text_frame.text = where
    style_paragraph(where_tb.text_frame.paragraphs[0],
                    size=9, italic=True, color=MID_GRAY)

    # Code block
    code_tb = s.shapes.add_textbox(Cm(x + 0.2), Cm(step_top + 1.4),
                                   Cm(step_w - 0.4), Cm(step_h - 1.6))
    code_tb.text_frame.word_wrap = True
    for j, line in enumerate(code.split("\n")):
        p = code_tb.text_frame.paragraphs[0] if j == 0 else code_tb.text_frame.add_paragraph()
        p.text = line if line else " "
        for r in p.runs:
            r.font.size = Pt(9)
            r.font.name = "Consolas"
            r.font.color.rgb = TEXT_DARK

    # Down arrow between boxes
    if i < n_steps - 1:
        ax = x + step_w + 0.02
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Cm(ax), Cm(step_top + step_h / 2 - 0.3),
            Cm(gap_w - 0.04), Cm(0.6),
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = MID_GRAY
        arr.line.fill.background()

# Bottom guarantees
guarantee = s.shapes.add_textbox(Cm(1.5), Cm(9.2), Cm(22.5), Cm(2.5))
write_paragraphs(
    guarantee.text_frame,
    [
        "Adapters never import each other — swapping one cannot break others.",
        "Pipeline calls only Protocol methods + the registry. Concrete adapters are "
        "instantiated by name; no hard-coded if/else over tool names.",
        "Config is the single dispatch knob. New tool = new adapter file + 1 line in "
        "config — no pipeline edit.",
    ],
    size=12, bullet=True, space_after=4,
)


# ---------- Slide 7: Output Contract ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Output Contract — Stable Across Tool Swaps")
remove_placeholder(s, 13)

# Left column: bullets (narrower so the diagram on the right has room)
left_box = s.shapes.add_textbox(Cm(1.0), Cm(3.2), Cm(10.0), Cm(9.0))
write_paragraphs(
    left_box.text_frame,
    [
        "Per-element JSON sidecars are the source of truth.",
        "content_list.json is a deterministic merge of those sidecars — "
        "rebuildable any time via assemble.",
        "Schema is stable: tool changes do not break downstream consumers.",
        "Run-fingerprint (PDF hash, DPI, stage config) prevents "
        "stale-output mix-ups.",
    ],
    size=14, bullet=True, space_after=10,
)

# Right column: element-types fan-out diagram (widened)
RX = 11.6      # right column origin (cm)
RW = 13.5      # right column width
CX = RX + RW / 2  # center x of right column
RY = 3.4

# Top row: text-only types
top_w = 4.6
top_gap = 0.6
top_total = 2 * top_w + top_gap
top_x0 = CX - top_total / 2
for i, lab in enumerate(["text", "heading"]):
    cx = top_x0 + i * (top_w + top_gap)
    labeled_box(s, cx, RY, top_w, 1.0, lab, fill=LIGHT_BG, size=12)
    # connector down to top of Element
    line = s.shapes.add_connector(
        1, Cm(cx + top_w / 2), Cm(RY + 1.0),
        Cm(CX), Cm(RY + 3.5),
    )
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

# Center "Element"
labeled_box(s, CX - 2.0, RY + 3.5, 4.0, 1.3,
            "Element", fill=BAM_RED, fg=WHITE,
            size=14, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)

# Bottom row: visual element types (wider boxes, smaller font)
bot_labels = ["table", "formula", "figure", "diagram", "tech. drawing"]
bot_w = 2.55
bot_gap = 0.10
bot_total = len(bot_labels) * bot_w + (len(bot_labels) - 1) * bot_gap
bot_x0 = CX - bot_total / 2
for i, lab in enumerate(bot_labels):
    cx = bot_x0 + i * (bot_w + bot_gap)
    labeled_box(s, cx, RY + 6.0, bot_w, 1.1, lab,
                fill=ACCENT_LIGHT, size=10)
    # connector up to bottom of Element
    line = s.shapes.add_connector(
        1, Cm(cx + bot_w / 2), Cm(RY + 6.0),
        Cm(CX), Cm(RY + 4.8),
    )
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

# Cluster labels
upper_lbl = s.shapes.add_textbox(Cm(RX), Cm(RY + 1.2), Cm(RW), Cm(0.6))
upper_lbl.text_frame.text = "text-only types  ·  JSON sidecar only"
style_paragraph(upper_lbl.text_frame.paragraphs[0],
                size=10, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

lower_lbl = s.shapes.add_textbox(Cm(RX), Cm(RY + 7.3), Cm(RW), Cm(0.6))
lower_lbl.text_frame.text = "visual types  ·  JSON sidecar + PNG crop"
style_paragraph(lower_lbl.text_frame.paragraphs[0],
                size=10, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)


# ---------- Slide 8: Current Status (v0.1.0) ----------
s = prs.slides.add_slide(prs.slide_layouts[6])  # 1-spaltig universal (no auto-numbering)
set_title(s, "Current Status  ·  v0.1.0")
body = get_placeholder(s, 13)
write_paragraphs(
    body.text_frame,
    [
        "End-to-end pipeline runs: segment → extract-text → "
        "describe-figures → assemble.",
        "Default GPU stack working: MinerU 3.1 segmenter + "
        "passthrough · Qwen2.5-VL figures.",
        "CPU-only fallback path available (pymupdf_text + noop) for "
        "smoke tests without GPU.",
        "Stage safety: PDF hash · DPI · stage-config fingerprinted "
        "in segmentation.json — stale outputs flagged, not overwritten.",
        "Per-element sidecars and content_list.json produced; "
        "assemble is byte-deterministic.",
        "Quality gates green: pytest · ruff · mypy.",
    ],
    size=14, bullet=True, space_after=6,
)


# ---------- Slide 9: Next Steps ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Next Steps")
remove_placeholder(s, 13)

# Left column: 5 audit tracks
audit_x, audit_y = 1.5, 3.0
audit_head = s.shapes.add_textbox(Cm(audit_x), Cm(audit_y), Cm(11.0), Cm(1.5))
audit_head.text_frame.word_wrap = True
audit_head.text_frame.text = "1.  Quality audits"
style_paragraph(audit_head.text_frame.paragraphs[0],
                size=18, bold=True, color=DARK_BLUE)
sub = audit_head.text_frame.add_paragraph()
sub.text = "extractor output vs. source PDFs"
style_paragraph(sub, size=12, italic=True, color=MID_GRAY)

audits = ["segmentation", "tables", "formulas", "figures", "merge"]
for i, name in enumerate(audits):
    cy = audit_y + 2.0 + i * 1.05
    labeled_box(s, audit_x + 0.5, cy, 10.0, 0.85, "·  " + name,
                fill=LIGHT_BG, size=13)

# Right column: Phase 2 outlook
ph2_x, ph2_y = 13.5, 3.0
ph2_head = s.shapes.add_textbox(Cm(ph2_x), Cm(ph2_y), Cm(11.0), Cm(1.5))
ph2_head.text_frame.word_wrap = True
ph2_head.text_frame.text = "2.  Phase 2 enrichment"
style_paragraph(ph2_head.text_frame.paragraphs[0],
                size=18, bold=True, color=DARK_BLUE)
sub2 = ph2_head.text_frame.add_paragraph()
sub2.text = "structural + relational layer above Phase 1"
style_paragraph(sub2, size=12, italic=True, color=MID_GRAY)

ph2_body = s.shapes.add_textbox(Cm(ph2_x), Cm(ph2_y + 2.0), Cm(11.0), Cm(7.5))
write_paragraphs(
    ph2_body.text_frame,
    [
        "Section hierarchy from headings / PDF outline.",
        "Mention parsing (“Tabelle 1”, “Abb. 3”, “(5)”).",
        "captioned_by + refers_to relations.",
        "document_rich.json artefact.",
        "All computable from Phase-1 output — no PDF re-parsing.",
    ],
    size=13, bullet=True, space_after=6,
)


# ---------- Slide 10: Discussion ----------
s = prs.slides.add_slide(prs.slide_layouts[6])
set_title(s, "Take-aways  ·  Discussion")
body = get_placeholder(s, 13)
write_paragraphs(
    body.text_frame,
    [
        "Modular extraction with a stable, tool-agnostic output contract.",
        "Tools are swappable per role — no pipeline rewrite to upgrade a model.",
        "Phase 1 is functional; quality audits and Phase-2 enrichment up next.",
        "",
        "Open questions for the room:",
        "•   Which reference PDFs should anchor the audits?",
        "•   Known extractor weaknesses from related projects?",
        "•   Phase 2 priorities: section tree, or mention/relation parsing?",
    ],
    size=14, bullet=False, space_after=4,
)
# Bullet styling: only the first three points get bullets manually
tf = body.text_frame
for i, p in enumerate(tf.paragraphs):
    if i < 3 and p.text:
        p.text = "• " + p.text
        style_paragraph(p, size=14, color=TEXT_DARK)
    elif p.text.startswith("Open questions"):
        style_paragraph(p, size=14, bold=True, color=DARK_BLUE)
    elif p.text.startswith("•"):
        style_paragraph(p, size=13, color=TEXT_DARK)


# --- Footer + slide numbers (skip the title slide) ----------------------
for slide in list(prs.slides)[1:]:
    add_footer_block(slide)


# --- Speaker notes ------------------------------------------------------
SPEAKER_NOTES = [
    # Slide 1 — Title
    "Welcome. Phase 1 of the AI-document-analysis project at FB 3.3. "
    "Short talk: vision, current state, what's next, then discussion.",

    # Slide 2 — Why & Goal
    "Two-level goal. Short term: extract every modality (text, tables, "
    "formulas, figures, drawings) into a uniform structure. Long term: "
    "verify technical claims and trace each one back to its supporting "
    "evidence. Concrete driver: BAM's approval of safety-critical "
    "transport containers, where claims need to be substantiated against "
    "the original source material.",

    # Slide 3 — Architecture
    "Four layers, each independently replaceable. Phase 1 is Layer 1, "
    "Extraction. The layers above (Embedding, Storage, Agent) consume "
    "Phase 1's output via a stable JSON contract. That decoupling is the "
    "whole point — we can later swap the extractor or any higher layer "
    "without breaking the others.",

    # Slide 4 — Pipeline
    "Four staged commands. Each runs in its own process so GPU memory is "
    "released between stages — important when shuffling MinerU's "
    "segmenter and Qwen's VLM in and out of GPU. The pipeline never "
    "decides which tool to use heuristically; YAML config is the only "
    "knob. Six swappable roles in total.",

    # Slide 5 — Code Architecture
    "Three concentric layers. Core: contracts — Protocols, registry, "
    "config, models, output. Stages: orchestration of the four staged "
    "commands. Adapters: concrete tool wrappers, lazy-imported so the "
    "GPU stack is only loaded if the active config asks for it. The CLI "
    "is the thin entry point on top — it just dispatches to a stage. "
    "Tests sit alongside the package; integration paths are marker-gated.",

    # Slide 6 — Modularity Guarantee
    "Four-step contract that keeps tool swaps safe. (1) A Protocol in "
    "interfaces.py defines what any adapter for that role must look "
    "like. (2) The adapter implements the protocol and self-registers "
    "via a decorator with a name. (3) Config picks the adapter by that "
    "name. (4) The pipeline only ever calls protocol methods through "
    "the registry — it never imports a concrete adapter directly. Net "
    "effect: adding a new tool means one new file plus one line in "
    "config; no pipeline edit, no chance of breaking other adapters.",

    # Slide 7 — Output Contract
    "Per-element JSON sidecars are the source of truth. content_list.json "
    "is just a deterministic merge — rebuildable any time via "
    "`python -m extraction assemble`. Schema is stable: as long as the "
    "config swap doesn't change the role names, downstream consumers "
    "don't notice. Run-fingerprint guards against stale outputs.",

    # Slide 8 — Status
    "Where we are today: pipeline runs end to end on real PDFs with the "
    "default GPU stack. CPU-only fallback path exists for smoke tests "
    "without GPU. Stage safety guards against running with stale outputs "
    "by fingerprinting PDF hash, render DPI, and the relevant config. "
    "Quality gates (pytest, ruff, mypy) are part of the standard test "
    "loop.",

    # Slide 9 — Next Steps
    "Two threads. (1) Quality audits — systematically compare extractor "
    "output against the source PDFs, per modality. Goal: surface "
    "weaknesses in segmentation, tables, formulas, figures, and the merge "
    "layer before Phase 2 ramps up. (2) Phase 2 will add structural "
    "enrichment — section tree, mention parsing, captioned-by / refers-to "
    "relations — fully computable from Phase 1 output without re-parsing "
    "the PDFs.",

    # Slide 10 — Discussion
    "Take-aways: stable contract, swappable tools, Phase 1 functional. "
    "Open questions for the room: which reference PDFs should anchor the "
    "audits? Which extractor weaknesses do you already know from related "
    "projects? Phase 2 priorities — section tree first, or "
    "mention/relation parsing first?",
]
for slide, notes in zip(prs.slides, SPEAKER_NOTES):
    set_speaker_notes(slide, notes)


# --- Save ---------------------------------------------------------------
DST.parent.mkdir(parents=True, exist_ok=True)
prs.save(DST)
print(f"Wrote {DST}")
